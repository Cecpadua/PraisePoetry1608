from pptx import Presentation
from PIL import Image, ImageEnhance
import pytesseract
from io import BytesIO
import cv2
import numpy as np

def enhance_text(image_path, contrast_factor=1.5, sharpness_factor=3.0):
    # Open the image
    img = Image.open(image_path)

    # Enhance contrast
    enhancer = ImageEnhance.Contrast(img)
    img = enhancer.enhance(contrast_factor)

    # Enhance sharpness
    enhancer = ImageEnhance.Sharpness(img)
    img = enhancer.enhance(sharpness_factor)

    # Save the enhanced image
    enhanced_path = "enhanced_image.png"
    img.save(enhanced_path)

    return enhanced_path

def preprocess_image(image_path):
    # Read the image
    img = cv2.imread(image_path)

    # Color filtering - keep only black
    lower_black = np.array([0, 0, 0], dtype=np.uint8)
    upper_black = np.array([50, 50, 50], dtype=np.uint8)
    mask = cv2.inRange(img, lower_black, upper_black)
    result = cv2.bitwise_and(img, img, mask=mask)

    # Convert to grayscale
    gray = cv2.cvtColor(result, cv2.COLOR_BGR2GRAY)

    # Binarization
    _, binary = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)

    # Save the preprocessed image
    cv2.imwrite("preprocessed_image.png", binary)

    return "preprocessed_image.png"

def extract_text_from_image(image_path):
    # Open the image and extract text using pytesseract
    img = Image.open(image_path)
    text = pytesseract.image_to_string(img, lang='chi_sim', config='--psm 6')
    return text

from pypinyin import pinyin, Style

def get_pinyin_initials(text):
    # 将中文字符串转换为拼音列表，Style.FIRST_LETTER 表示获取首字母
    pinyin_list = pinyin(text, style=Style.FIRST_LETTER)
    
    # 将拼音列表转换为字符串
    initials = ''.join([item[0] for item in pinyin_list])
    
    return initials


from cnocr import CnOcr
def extract_text_from_image2(image_path):
    ocr = CnOcr(det_model_name='naive_det')
    out = ocr.ocr(image_path)
    return '\n'.join([x['text'] for x in out])

def set_all_recognized_text_as_comments(pptx_path, filename):
    presentation = Presentation(pptx_path)
    all_recognized_text = ""
    all_recognized_py = ""

    for slide_number, slide in enumerate(presentation.slides, start=1):
        print(f'Processing Slide {slide_number}...')
        for shape_number, shape in enumerate(slide.shapes, start=1):
            print(f'Processing Shape {shape_number}...')
            # Check if the shape has an image
            if shape.shape_type == 13:  # 13 corresponds to Picture shape type
                print(f'Extracting text from Image in Shape {shape_number}...')
                image_path = f'image_slide{slide_number}_shape{shape_number}.png'

                # Save the image using BytesIO
                img_bytes = BytesIO(shape.image.blob)
                img = Image.open(img_bytes)
                img = img.convert('RGB')
                img.save(image_path)

                # Preprocess the image
                preprocessed_image_path = preprocess_image(image_path)
                enhance_text_image_path = enhance_text(preprocessed_image_path)


                text = extract_text_from_image2(enhance_text_image_path).strip().replace(' ', '').replace('.', '')
                text = text.replace("(", '').replace(")", '').replace("（", '').replace("）", '').replace(":", '')\
                    .replace("“", '').replace("”", '').replace("‘", '').replace("’", '').replace(";", '').replace("；", '')\
                    .replace("。", '').replace(",", '').replace("、", '').replace("？", '').replace("！", '').replace("!", '')\
                    .replace("?", '').replace("_", '').replace("C", '').replace("D", '').replace("E", '').replace("F", '')\
                    .replace("G", '').replace("A", '').replace("B", '')

                # Split into lines
                lines = text.split('\n')
                # Remove lines with fewer than 5 characters
                filtered_lines = [line for line in lines if len(line) >= 5]
                # Reassemble into new text
                filtered_text = '\n'.join(filtered_lines)
                
                # Print the result
                # print(f'[Result2]: {filtered_text}')
                initials = get_pinyin_initials(filtered_text)
                # print(f'[Initials]: {initials}')
                
                # Append to the overall recognized text
                all_recognized_text += f'\n{filtered_text}'
                all_recognized_py += f'\n{initials}'
                # print("-------------------------------------------------")
    # Set all recognized text as comments in the core_properties
    #print(all_recognized_text)
    presentation.core_properties.title = filename.replace('.pptx', '')
    presentation.core_properties.comments = all_recognized_text[0:255]
    presentation.core_properties.subject = all_recognized_py[0:255]
    presentation.core_properties.version = "1.1"
    presentation.core_properties.language = "zh-CN"
    presentation.core_properties.keywords = "投影"
    presentation.core_properties.category = "投影"
    presentation.core_properties.content_status = "索引完成"
    presentation.core_properties.author = "Yihao Zhuo"

    # Save the modified PPTX file
    presentation.save(pptx_path)

    # Save the extracted text to a file
    with open("D:/Users/yihao zhuo/Desktop/ppt/dir/ppt 16-9/" + filename + ".txt", 'w', encoding='utf-8') as file:
        file.write(all_recognized_text)
        file.write(all_recognized_py)
        file.close()

# Usage example
# pptx_path = 'D:/Users/yihao zhuo/Desktop/ppt/dir/test/0002.亚伯拉罕之神.pptx'
# set_all_recognized_text_as_comments(pptx_path, "0002.亚伯拉罕之神")

import os
import re
def extract_numeric_part(file_name):
    numeric_part = re.search(r'\d+', file_name)
    if numeric_part:
        return int(numeric_part.group())
    else:
        return None
    
for file in os.listdir("D:/Users/yihao zhuo/Desktop/ppt/dir/ppt 16-9"):
    print(file)
    startIndex = 614
    if file.endswith(".pptx") and int(extract_numeric_part(file)) >= startIndex:
        pptx_path = f"D:/Users/yihao zhuo/Desktop/ppt/dir/ppt 16-9/{file}"
        set_all_recognized_text_as_comments(pptx_path, file.replace('.pptx', ''))
        print(f"Processed {file}")
